#!/usr/bin/env python3
"""
Extract content from .pptx files for conversion to web presentations.

Usage:
    python extract-pptx.py <input.pptx> [output_dir]

Requires: pip install python-pptx

Adapted from frontend-slides by zarazhangrui.
"""

import json
import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)


def extract_pptx(file_path, output_dir=None):
    """Extract slides, text, images, and notes from a .pptx file."""
    if not os.path.exists(file_path):
        print(f"Error: File not found: {file_path}")
        sys.exit(1)

    if output_dir is None:
        output_dir = os.path.splitext(file_path)[0] + "_extracted"

    assets_dir = os.path.join(output_dir, "assets")
    os.makedirs(assets_dir, exist_ok=True)

    prs = Presentation(file_path)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": slide_num,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
        }

        for shape in slide.shapes:
            # Extract title
            if shape.has_text_frame and shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                slide_info["title"] = shape.text_frame.text.strip()
                continue

            # Extract text content
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_info["content"].append(text)

            # Extract images (shape type 13 = Picture)
            if shape.shape_type == 13:
                image = shape.image
                image_bytes = image.blob
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"

                image_filename = f"slide{slide_num}_img{len(slide_info['images']) + 1}.{ext}"
                image_path = os.path.join(assets_dir, image_filename)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                slide_info["images"].append({
                    "filename": image_filename,
                    "path": image_path,
                    "width": shape.width,
                    "height": shape.height,
                })

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                slide_info["notes"] = notes_frame.text.strip()

        slides_data.append(slide_info)

    # Write JSON output
    output_path = os.path.join(output_dir, "extracted-slides.json")
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(slides_data, f, ensure_ascii=False, indent=2)

    # Print summary
    print(f"\nExtraction complete:")
    print(f"  Slides: {len(slides_data)}")
    print(f"  Output: {output_path}")
    print(f"  Assets: {assets_dir}")
    print()
    for s in slides_data:
        img_count = len(s["images"])
        img_info = f" ({img_count} image{'s' if img_count != 1 else ''})" if img_count else ""
        title = s["title"] or "(no title)"
        print(f"  Slide {s['slide_number']}: {title}{img_info}")

    return slides_data


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(f"Usage: python {sys.argv[0]} <input.pptx> [output_dir]")
        sys.exit(1)

    input_file = sys.argv[1]
    out_dir = sys.argv[2] if len(sys.argv) > 2 else None
    extract_pptx(input_file, out_dir)
